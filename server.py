from fastmcp import FastMCP
import re
import os
from pptx import Presentation
from typing import Dict, List, Any, Optional
from config import get_template_path, get_etapa_regex, get_placeholders, validate_template_exists

mcp = FastMCP("mcp-mapeamento")

def verificar_arquivo_em_uso(caminho_arquivo: str) -> bool:
    """
    Verifica se um arquivo está sendo usado por outro processo.
    
    Args:
        caminho_arquivo: Caminho completo do arquivo
        
    Returns:
        True se o arquivo está em uso, False caso contrário
    """
    try:
        # Tentar abrir o arquivo em modo de escrita exclusiva
        with open(caminho_arquivo, 'a'):
            return False
    except (IOError, OSError):
        return True

def gerar_nome_arquivo_unico(diretorio: str, nome_base: str, extensao: str = ".pptx") -> str:
    """
    Gera um nome de arquivo único no diretório especificado.
    
    Args:
        diretorio: Diretório onde o arquivo será salvo
        nome_base: Nome base do arquivo (sem extensão)
        extensao: Extensão do arquivo (padrão: .pptx)
        
    Returns:
        Nome de arquivo único
    """
    nome_arquivo = f"{nome_base}{extensao}"
    caminho_completo = os.path.join(diretorio, nome_arquivo)
    
    contador = 1
    while os.path.exists(caminho_completo) or verificar_arquivo_em_uso(caminho_completo):
        nome_arquivo = f"{nome_base}_{contador}{extensao}"
        caminho_completo = os.path.join(diretorio, nome_arquivo)
        contador += 1
    
    return nome_arquivo

def processar_etapas(texto: str) -> List[Dict[str, str]]:
    """
    Processa o texto das etapas e extrai as informações necessárias.
    
    Args:
        texto: Texto contendo as etapas no formato especificado
        
    Returns:
        Lista de dicionários com as informações de cada etapa
    """
    etapas = []
    
    # Regex para capturar cada etapa
    pattern = get_etapa_regex()
    
    matches = re.findall(pattern, texto, re.DOTALL)
    
    for match in matches:
        numero, nome, supplier, input_info, process, output, customer = match
        
        # Criar descrição baseada no processo
        descricao = f"Esta etapa envolve {process.lower()}"
        
        # Criar regras baseadas em input e output
        regras = f"Input: {input_info.strip()}\nOutput: {output.strip()}\nSupplier: {supplier.strip()}\nCustomer: {customer.strip()}"
        
        etapas.append({
            'numero': numero.strip(),
            'nome': nome.strip(),
            'descricao': descricao,
            'regras': regras
        })
    
    return etapas

def criar_apresentacao_do_template(template_path: str = None) -> Presentation:
    """
    Cria uma nova apresentação baseada no template.
    
    Args:
        template_path: Caminho para o template PPTX (opcional, usa config se não fornecido)
        
    Returns:
        Objeto Presentation
    """
    if template_path is None:
        template_path = get_template_path()
    
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template não encontrado: {template_path}")
    
    return Presentation(template_path)

@mcp.tool()
def criar_ddp(
    texto_etapas: str,
    nome_processo: str,
    diretorio_saida: str,
    nome_arquivo: Optional[str] = None
) -> Dict[str, Any]:

    """
    Esta ferramenta ("criar_ddp") gera automaticamente uma apresentação PowerPoint para o "Documento de Desenho de Processo" (DDP) a partir de um texto estruturado fornecido pelo usuário contendo as etapas de um processo. Ela utiliza um template pré-definido e preenche os slides com as informações extraídas do texto, organizando cada etapa em um slide separado, incluindo nome, descrição e regras de cada etapa.

    Exemplos de chamada:
    "Quero montar um novo DDP"
    "Quero montar um novo DDP para o processo"
    "Montar um novo DDP para o processo"
    Entre outros parecidos.

    Caso o usuário não forneça diretamente os parâmetros abaixo, solicite para ele.

    Parâmetros obrigatórios:
    - texto_etapas: Um texto contendo as etapas do processo, seguindo o formato esperado (com campos como número da etapa, nome, supplier, input, process, output, customer).
    - nome_processo: O nome do processo que está sendo mapeado, que será utilizado no título e no nome do arquivo.
    - diretorio_saida: O diretório onde o arquivo PPTX gerado será salvo.
    - nome_arquivo: (Opcional) Nome do arquivo de saída. Se não for fornecido, será gerado automaticamente no formato DDP_{nome_processo}.pptx.

    O retorno será um dicionário informando o sucesso, o caminho do arquivo gerado, o nome do processo, o total de etapas e detalhes dos slides criados.
    """
    etapas = processar_etapas(texto_etapas)
    if not etapas:
        return {"erro": "Nenhuma etapa válida encontrada no texto fornecido. Verifique o formato."}

    presentation = criar_apresentacao_do_template()
    slides_criados = []

    for etapa in etapas:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        # Substitui os placeholders pelo nome, igual ao exemplo do debug_local.py
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                if shape.name == 'Text Placeholder 2':
                    shape.text_frame.text = f"ETAPA {etapa['numero']} - {etapa['nome']}"
                elif shape.name == 'Text Placeholder 3':
                    shape.text_frame.text = etapa['descricao']
                elif shape.name == 'Text Placeholder 1':
                    shape.text_frame.text = etapa['regras']
        slides_criados.append({
            'slide_index': len(presentation.slides) - 1,
            'etapa': f"ETAPA {etapa['numero']} - {etapa['nome']}"
        })

    os.makedirs(diretorio_saida, exist_ok=True)

    if nome_arquivo is None:
        nome_base = f"DDP_{nome_processo.replace(' ', '_')}"
        nome_arquivo = gerar_nome_arquivo_unico(diretorio_saida, nome_base)
    else:
        nome_sem_extensao = os.path.splitext(nome_arquivo)[0]
        extensao = os.path.splitext(nome_arquivo)[1] or ".pptx"
        nome_arquivo = gerar_nome_arquivo_unico(diretorio_saida, nome_sem_extensao, extensao)

    caminho_completo = os.path.join(diretorio_saida, nome_arquivo)
    presentation.save(caminho_completo)

    return {
        "sucesso": True,
        "mensagem": "Apresentação DDP criada com sucesso!",
        "arquivo": caminho_completo,
        "nome_processo": nome_processo,
        "total_etapas": len(etapas),
        "slides_criados": slides_criados,
        "etapas_processadas": etapas
    }


if __name__ == "__main__":
    mcp.run()